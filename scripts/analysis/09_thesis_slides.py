"""
Generates a 15-slide PowerPoint presentation from thesis_minimap2_gradient.md.

Slide structure:
  1.  Title
  2.  The Problem
  3.  Research Objectives
  4.  Biological Background — PAIs and HGT
  5.  Prior Work and Methodological Gap
  6.  Pipeline Overview (workflow diagram as text)
  7.  minimap2 Divergence Gradient — How It Works
  8.  Marker Score — Ranking the 415 Candidates
  9.  Threshold Sensitivity (table)
  10. Chromosome Divergence Landscape (figure)
  11. Confidence Tiers — Intersection Filter (figure)
  12. ML Classifier — Features and Performance
  13. Ablation Study — Feature Contributions
  14. Negative Control — Specificity Validation
  15. Conclusion and Future Directions

Output: data/results/figures/thesis_presentation.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import textwrap

REPO_ROOT = Path(__file__).resolve().parents[2]
FIG_DIR   = REPO_ROOT / "data" / "results" / "figures"
OUT_PATH  = FIG_DIR / "thesis_presentation.pptx"

# ── Colour palette ─────────────────────────────────────────────────────────────
C_NAVY    = RGBColor(0x1A, 0x2B, 0x4A)   # slide background / title bar
C_TEAL    = RGBColor(0x00, 0x7A, 0x8A)   # accent / subheading
C_AMBER   = RGBColor(0xF3, 0x9C, 0x12)   # highlight
C_RED     = RGBColor(0xE7, 0x4C, 0x3C)   # HIGH-gradient
C_GREEN   = RGBColor(0x27, 0xAE, 0x60)   # positive / Tier 1
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xF4, 0xF6, 0xF8)   # body background
C_DARK    = RGBColor(0x2C, 0x3E, 0x50)   # body text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color: RGBColor,
             line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=Pt(18), bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txbox


def add_para(tf, text, font_size=Pt(16), bold=False,
             color=C_DARK, align=PP_ALIGN.LEFT, space_before=Pt(4)):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def header_bar(slide, title: str, subtitle: str = ""):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), C_NAVY)
    add_text(slide, title,
             Inches(0.35), Inches(0.1), Inches(12), Inches(0.6),
             font_size=Pt(26), bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.72), Inches(12), Inches(0.38),
                 font_size=Pt(14), color=C_TEAL)


def footer(slide, slide_num: int, total: int = 15):
    add_rect(slide, 0, SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32), C_NAVY)
    add_text(slide,
             "Oluwatosin Samuel Oluwole | MSc Microbiology | Carl von Ossietzky University | May 2026",
             Inches(0.3), SLIDE_H - Inches(0.32), Inches(10), Inches(0.32),
             font_size=Pt(9), color=C_WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, f"{slide_num} / {total}",
             Inches(12.5), SLIDE_H - Inches(0.32), Inches(0.8), Inches(0.32),
             font_size=Pt(9), color=C_WHITE, align=PP_ALIGN.RIGHT)


def add_image(slide, path: Path, left, top, width, height=None):
    if not path.exists():
        return
    if height:
        slide.shapes.add_picture(str(path), left, top, width, height)
    else:
        slide.shapes.add_picture(str(path), left, top, width)


def bullet_frame(slide, left, top, width, height, bg=C_LIGHT):
    add_rect(slide, left, top, width, height, bg)
    txbox = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.1),
                                      width - Inches(0.25), height - Inches(0.15))
    tf = txbox.text_frame
    tf.word_wrap = True
    return tf


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_NAVY)

    # Central white card
    add_rect(sl, Inches(0.7), Inches(1.2), Inches(11.9), Inches(4.5), C_WHITE)

    add_text(sl,
             "Hybrid Pangenomic and Sequence-Based Framework\n"
             "for Reducing False Positives in Pathogen Detection\n"
             "from Metagenomic Data",
             Inches(1.0), Inches(1.4), Inches(11.3), Inches(2.2),
             font_size=Pt(28), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    add_text(sl,
             "minimap2 Divergence Gradient Methodology",
             Inches(1.0), Inches(3.55), Inches(11.3), Inches(0.45),
             font_size=Pt(17), color=C_TEAL, align=PP_ALIGN.CENTER)

    add_text(sl,
             "Oluwatosin Samuel Oluwole",
             Inches(1.0), Inches(4.05), Inches(11.3), Inches(0.38),
             font_size=Pt(15), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    add_text(sl,
             "MSc Microbiology  ·  Carl von Ossietzky University  ·  Supervisor: Prof. Denis Shields  ·  May 2026",
             Inches(1.0), Inches(4.42), Inches(11.3), Inches(0.38),
             font_size=Pt(12), color=C_DARK, align=PP_ALIGN.CENTER)

    # Bottom tag strip
    for i, (label, col) in enumerate([
        ("Comparative Genomics", C_TEAL),
        ("Machine Learning", C_AMBER),
        ("E. coli O157:H7", C_RED),
        ("Pathogen Discovery", C_GREEN),
    ]):
        x = Inches(1.4 + i * 2.65)
        add_rect(sl, x, Inches(5.85), Inches(2.35), Inches(0.45), col)
        add_text(sl, label, x, Inches(5.85), Inches(2.35), Inches(0.45),
                 font_size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    footer(sl, 1)
    return sl


def slide_02_problem(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_LIGHT)
    header_bar(sl, "The Problem", "Why do metagenomic classifiers generate false positives?")

    tf = bullet_frame(sl, Inches(0.4), Inches(1.3), Inches(6.0), Inches(5.5))
    bullets = [
        ("The core challenge", True, C_NAVY, Pt(17)),
        ("Metagenomic classifiers detect pathogens by aligning short reads\n"
         "to reference genomes. But E. coli O157:H7 and harmless commensal\n"
         "E. coli share >85% of their genomic sequence.", False, C_DARK, Pt(14)),
        ("", False, C_DARK, Pt(6)),
        ("What goes wrong", True, C_NAVY, Pt(17)),
        ("Commensal reads map to conserved regions of the pathogen reference\n"
         "→ the classifier calls a pathogen even when none is present.\n"
         "This is a false positive caused by phylogenetic relatedness, not\n"
         "a tool deficiency.", False, C_DARK, Pt(14)),
        ("", False, C_DARK, Pt(6)),
        ("The solution framing", True, C_NAVY, Pt(17)),
        ("Restrict classification to regions that DIFFER between pathogen\n"
         "and commensal — pathogenicity islands (PAIs). Extract a quantitative\n"
         "divergence signal from these regions and use it as an ML feature.", False, C_DARK, Pt(14)),
    ]
    first = True
    for text, bold, color, size in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.size = size

    # Right panel — stat callouts
    for top, val, label, col in [
        (Inches(1.35), ">85%", "Shared genome between\nE. coli O157:H7\nand commensals", C_TEAL),
        (Inches(3.2),  "415",  "Candidate pathogen-\nspecific marker regions\nidentified", C_AMBER),
        (Inches(5.05), "10.9×", "Enrichment of HIGH-gradient\nwindows: pathogen vs.\ncommensal comparison", C_RED),
    ]:
        add_rect(sl, Inches(6.8), top, Inches(5.85), Inches(1.6), col)
        add_text(sl, val,
                 Inches(6.8), top + Inches(0.1), Inches(5.85), Inches(0.75),
                 font_size=Pt(36), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, label,
                 Inches(6.8), top + Inches(0.8), Inches(5.85), Inches(0.7),
                 font_size=Pt(12), color=C_WHITE, align=PP_ALIGN.CENTER)

    footer(sl, 2)
    return sl


def slide_03_objectives(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_LIGHT)
    header_bar(sl, "Research Objectives", "What this thesis sets out to do")

    objectives = [
        ("1", "Develop a quantitative divergence gradient",
         "Represent pathogen-vs-commensal nucleotide divergence as structured,\n"
         "window-level features derived from minimap2 pairwise alignments\n"
         "across a 30-genome commensal panel.", C_TEAL),
        ("2", "Rank candidate marker regions",
         "Identify the 415 regions of the Sakai genome most consistently\n"
         "absent from commensals, ranked by a composite marker score\n"
         "(divergence × flank conservation × log-length).", C_AMBER),
        ("3", "Evaluate ML discriminative power",
         "Test whether adding gradient features to an XGBoost classifier\n"
         "improves AUROC over the established 10-feature NUCmer baseline,\n"
         "under leakage-free cross-validation.", C_RED),
        ("4", "Validate specificity",
         "Confirm that HIGH-gradient regions are genuinely pathogen-specific\n"
         "using pathogen-vs-pathogen negative controls and known PAI\n"
         "coordinate overlap.", C_GREEN),
    ]

    for i, (num, title, body, col) in enumerate(objectives):
        row = i // 2
        col_n = i % 2
        left = Inches(0.4 + col_n * 6.45)
        top  = Inches(1.35 + row * 2.85)
        add_rect(sl, left, top, Inches(6.1), Inches(2.65), col)
        add_text(sl, num,
                 left + Inches(0.12), top + Inches(0.08), Inches(0.5), Inches(0.6),
                 font_size=Pt(30), bold=True, color=C_WHITE)
        add_text(sl, title,
                 left + Inches(0.6), top + Inches(0.1), Inches(5.3), Inches(0.45),
                 font_size=Pt(15), bold=True, color=C_WHITE)
        add_text(sl, body,
                 left + Inches(0.15), top + Inches(0.6), Inches(5.8), Inches(1.95),
                 font_size=Pt(12), color=C_WHITE)

    footer(sl, 3)
    return sl


def slide_04_biology(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_LIGHT)
    header_bar(sl, "Biological Background",
               "Pathogenicity islands and horizontal gene transfer")

    tf = bullet_frame(sl, Inches(0.4), Inches(1.3), Inches(7.5), Inches(5.6))
    content = [
        ("What is a Pathogenicity Island (PAI)?", True, C_NAVY, Pt(16)),
        ("Large genomic regions (typically >10 kb) present in pathogenic\n"
         "strains and absent from non-pathogenic relatives of the same species\n"
         "(Hacker & Kaper 2000, Annu Rev Microbiol).", False, C_DARK, Pt(13)),
        ("", False, C_DARK, Pt(5)),
        ("Structural hallmarks", True, C_TEAL, Pt(15)),
        ("• Often flanked by direct repeats and tRNA genes\n"
         "• Carry mobility genes (integrase, transposase)\n"
         "• Deviant GC% and codon usage from host genome\n"
         "• Gradually ameliorate to host composition over time\n"
         "  (Lawrence & Ochman 1997, J Mol Evol)", False, C_DARK, Pt(13)),
        ("", False, C_DARK, Pt(5)),
        ("Key PAIs in E. coli O157:H7 Sakai", True, C_TEAL, Pt(15)),
        ("• LEE (Locus of Enterocyte Effacement) — type III secretion\n"
         "• OI-1, OI-7, OI-36, OI-43/48, OI-57 — virulence gene clusters\n"
         "• Sp1, Sp4, Sp11, Sp15 — Sakai prophage islands\n"
         "• pO157 megaplasmid — additional virulence genes", False, C_DARK, Pt(13)),
        ("", False, C_DARK, Pt(5)),
        ("Why divergence signals PAI regions", True, C_RED, Pt(15)),
        ("A window of the pathogen genome that fails to align\n"
         "to ANY commensal genome is likely pathogen-specific.\n"
         "divergence_score = 1 − (coverage × identity)", False, C_DARK, Pt(13)),
    ]
    first = True
    for text, bold, color, size in content:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.size = size

    # PAI location diagram — colour-coded boxes
    add_text(sl, "Known PAI locations on Sakai chromosome (5.5 Mb)",
             Inches(8.1), Inches(1.35), Inches(4.9), Inches(0.35),
             font_size=Pt(11), bold=True, color=C_NAVY)

    pai_data = [
        ("LEE",       2.796, "#c0392b"),
        ("OI-1",      0.043, "#8e44ad"),
        ("OI-7",      0.539, "#8e44ad"),
        ("OI-36",     1.667, "#8e44ad"),
        ("OI-43/48",  2.148, "#d35400"),
        ("OI-57",     2.480, "#8e44ad"),
        ("Sp1",       1.335, "#27ae60"),
        ("Sp4",       2.585, "#27ae60"),
        ("Sp11",      4.530, "#27ae60"),
        ("Sp15",      5.130, "#27ae60"),
    ]
    chrom_left  = Inches(8.1)
    chrom_top   = Inches(1.85)
    chrom_w     = Inches(4.85)
    chrom_h     = Inches(0.22)
    chrom_len   = 5.5

    add_rect(sl, chrom_left, chrom_top, chrom_w, chrom_h, RGBColor(0xBD, 0xC3, 0xC7))
    for name, pos_mb, hex_col in pai_data:
        x_frac = pos_mb / chrom_len
        px = chrom_left + Emu(int(chrom_w * x_frac))
        r, g, b = int(hex_col[1:3], 16), int(hex_col[3:5], 16), int(hex_col[5:7], 16)
        add_rect(sl, px, chrom_top, Inches(0.08), chrom_h, RGBColor(r, g, b))

    # Legend
    for i, (label, hex_col) in enumerate([
        ("LEE (T3SS)", "#c0392b"),
        ("OI islands", "#8e44ad"),
        ("OI-43/48",   "#d35400"),
        ("Prophage Sp", "#27ae60"),
    ]):
        r, g, b = int(hex_col[1:3], 16), int(hex_col[3:5], 16), int(hex_col[5:7], 16)
        top_l = Inches(2.2 + i * 0.38)
        add_rect(sl, Inches(8.1), top_l, Inches(0.22), Inches(0.25), RGBColor(r, g, b))
        add_text(sl, label, Inches(8.4), top_l, Inches(4.5), Inches(0.28),
                 font_size=Pt(11), color=C_DARK)

    footer(sl, 4)
    return sl


def slide_05_gap(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_LIGHT)
    header_bar(sl, "Prior Work and Methodological Gap",
               "What existed and what was missing")

    # Two-column comparison
    for col_n, (title, col, items) in enumerate([
        ("NUCmer approach (thesis-enhanced branch)", C_TEAL, [
            "Aligns commensal genomes to Sakai using NUCmer",
            "Classifies alignment blocks as CONSERVED / MODERATE / DIVERGED",
            "Per-block identity threshold only",
            "Blocks vary in length — hard to aggregate across comparisons",
            "Single threshold, no gradient resolution",
            "AUROC 0.717 (10-feature baseline)",
        ]),
        ("minimap2 Gradient approach (this branch)", C_AMBER, [
            "Aligns 30 commensals to Sakai using minimap2 -x asm5",
            "Fixed 500 bp windows tiled across entire genome",
            "Combined score: 1 − (coverage × identity)",
            "Mean aggregated across all 30 comparisons per window",
            "Three-state gradient: LOW / MID / HIGH",
            "Flanking conservation identifies insertion context",
        ]),
    ]):
        left = Inches(0.4 + col_n * 6.45)
        add_rect(sl, left, Inches(1.3), Inches(6.1), Inches(0.45), col)
        add_text(sl, title, left + Inches(0.1), Inches(1.3), Inches(5.9), Inches(0.45),
                 font_size=Pt(12), bold=True, color=C_WHITE)
        tf = bullet_frame(sl, left, Inches(1.78), Inches(6.1), Inches(4.35), C_WHITE)
        first = True
        for item in items:
            if first:
                p = tf.paragraphs[0]; first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"• {item}"
            run.font.size = Pt(13)
            run.font.color.rgb = C_DARK

    add_rect(sl, Inches(0.4), Inches(6.25), Inches(12.53), Inches(0.58),
             RGBColor(0xEB, 0xF5, 0xFB))
    add_text(sl,
             "Key finding: The gradient features carry biological signal (AUROC 0.549 standalone) "
             "but do not improve over the NUCmer baseline (0.709 combined vs 0.717 NUCmer). "
             "Contribution: a new operational representation, not a new classifier.",
             Inches(0.55), Inches(6.28), Inches(12.2), Inches(0.52),
             font_size=Pt(12), color=C_NAVY, align=PP_ALIGN.CENTER)

    footer(sl, 5)
    return sl


def slide_06_pipeline(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_LIGHT)
    header_bar(sl, "Pipeline Overview", "From raw genomes to ML-validated pathogen markers")

    steps = [
        ("s01", "Genome\nIngestion", "60 genomes\n(30 commensal\n+ 30 control)", C_TEAL),
        ("s02b", "minimap2\nAlignment", "asm5 preset\n500 bp windows\n30-commensal panel", C_AMBER),
        ("s03", "NUCmer\nMarkers", "149 unique\nalignment blocks\n(validation layer)", C_TEAL),
        ("s04", "Metagenome\nSimulation", "Simulated reads\nfor ML evaluation", C_AMBER),
        ("s07", "XGBoost\nClassifier", "15 features\nAUROC 0.709\nStratGroupKFold", C_RED),
        ("s08", "Biological\nValidation", "PAI overlap\nK-12 absence\nNeg. controls", C_GREEN),
    ]

    box_w = Inches(1.9)
    box_h = Inches(1.8)
    gap   = Inches(0.22)
    top_y = Inches(2.0)

    for i, (sid, title, body, col) in enumerate(steps):
        x = Inches(0.35) + i * (box_w + gap)
        add_rect(sl, x, top_y, box_w, box_h, col)
        add_text(sl, sid, x, top_y + Inches(0.05), box_w, Inches(0.3),
                 font_size=Pt(9), color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, title, x, top_y + Inches(0.28), box_w, Inches(0.6),
                 font_size=Pt(13), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, body, x, top_y + Inches(0.92), box_w, Inches(0.85),
                 font_size=Pt(10), color=C_WHITE, align=PP_ALIGN.CENTER)
        # Arrow (except last)
        if i < len(steps) - 1:
            ax = x + box_w + Inches(0.04)
            add_text(sl, "→", ax, top_y + Inches(0.7), gap, Inches(0.4),
                     font_size=Pt(18), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    # Output callouts
    callouts = [
        (Inches(0.35),  "window_scores_500bp.tsv"),
        (Inches(4.7),   "candidates_schemeA_500bp.tsv"),
        (Inches(9.05),  "feature_matrix.tsv"),
    ]
    for cx, label in callouts:
        add_rect(sl, cx, Inches(4.1), Inches(3.5), Inches(0.35),
                 RGBColor(0xD5, 0xD8, 0xDC))
        add_text(sl, label, cx + Inches(0.08), Inches(4.1), Inches(3.3), Inches(0.35),
                 font_size=Pt(10), color=C_DARK)

    add_text(sl,
             "Negative controls: pathogen-vs-pathogen panel (5 EHEC strains) + label shuffle\n"
             "Confidence tiers: Tier 1 = both NUCmer + minimap2 agree (333 markers) | Tier 2 = NUCmer-only (82 markers)",
             Inches(0.4), Inches(4.65), Inches(12.5), Inches(0.65),
             font_size=Pt(12), color=C_NAVY, align=PP_ALIGN.CENTER)

    footer(sl, 6)
    return sl


def slide_07_gradient(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_LIGHT)
    header_bar(sl, "minimap2 Divergence Gradient — How It Works",
               "Converting alignments into a window-level divergence feature")

    # Left: formula steps
    steps = [
        ("Step 1 — Align each commensal to Sakai",
         "minimap2 -x asm5 commensal.fasta sakai.fasta → PAF file\n"
         "(PAF = Pairwise mApping Format — one alignment block per line)", C_TEAL),
        ("Step 2 — Score each 500 bp window",
         "For each aligned block overlapping the window:\n"
         "  divergence_score = 1 − (coverage_fraction × identity_fraction)\n"
         "  → 0 means perfectly conserved, 1 means completely absent", C_AMBER),
        ("Step 3 — Aggregate across 30 commensals",
         "Take the mean divergence score across all 30 commensal comparisons.\n"
         "A window divergent from ALL commensals is a stronger signal.", C_RED),
        ("Step 4 — Classify into gradient states",
         "LOW  ≤ 0.20  (conserved — skip for PAI detection)\n"
         "MID    0.20–0.60  (partial divergence)\n"
         "HIGH > 0.60  (candidate pathogen-specific region)", C_GREEN),
    ]

    for i, (title, body, col) in enumerate(steps):
        top = Inches(1.35 + i * 1.46)
        add_rect(sl, Inches(0.4), top, Inches(0.06), Inches(1.3), col)
        add_text(sl, title, Inches(0.55), top, Inches(6.5), Inches(0.38),
                 font_size=Pt(13), bold=True, color=col)
        add_text(sl, body, Inches(0.55), top + Inches(0.38), Inches(6.5), Inches(0.9),
                 font_size=Pt(12), color=C_DARK)

    # Right panel — gradient visual
    add_text(sl, "Gradient classification (Scheme A)",
             Inches(7.5), Inches(1.35), Inches(5.45), Inches(0.38),
             font_size=Pt(13), bold=True, color=C_NAVY)

    for top, label, val_range, col in [
        (Inches(1.85), "LOW  ≤ 0.20",    "Shared with commensals\nConserved housekeeping genes", C_TEAL),
        (Inches(3.1),  "MID  0.20–0.60", "Partial divergence\nRecombining or partially transferred regions", C_AMBER),
        (Inches(4.35), "HIGH > 0.60",    "Consistently absent from commensals\n→ Candidate PAI / pathogen-specific region", C_RED),
    ]:
        add_rect(sl, Inches(7.5), top, Inches(5.45), Inches(1.0), col)
        add_text(sl, label, Inches(7.65), top + Inches(0.05), Inches(5.2), Inches(0.38),
                 font_size=Pt(14), bold=True, color=C_WHITE)
        add_text(sl, val_range, Inches(7.65), top + Inches(0.42), Inches(5.2), Inches(0.5),
                 font_size=Pt(11), color=C_WHITE)

    add_text(sl,
             "10,889 windows tiled across 5.5 Mb Sakai chromosome  ·  30 commensal comparisons per window\n"
             "2,870 HIGH windows (25.65%) merged into 143 candidate regions under Scheme A",
             Inches(7.5), Inches(5.55), Inches(5.45), Inches(0.7),
             font_size=Pt(11), color=C_NAVY)

    footer(sl, 7)
    return sl


def slide_08_marker_score(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_LIGHT)
    header_bar(sl, "Marker Score — Ranking the 415 Candidate Regions",
               "A composite statistic combining divergence, context, and size")

    add_text(sl, "marker_score  =  mean_divergence  ×  flank_conservation  ×  log₁₊(region_length)",
             Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.55),
             font_size=Pt(18), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(0.4), Inches(1.88), Inches(12.5), Inches(0.04), C_TEAL)

    terms = [
        ("mean_divergence", "0 – 1",
         "Average divergence score across all\n30 commensal comparisons for this region.\n"
         "Core signal of pathogen specificity.",
         C_RED),
        ("flank_conservation", "0 – 1",
         "1 − mean divergence of the 2 kb flanking\nsequence on either side of the candidate.\n"
         "High value = conserved neighbourhood\n= structural signature of HGT insertion.",
         C_TEAL),
        ("log₁₊(region_length)", "> 0",
         "Length bonus with diminishing returns.\n"
         "6.5 kb block → 8.78   vs   500 bp → 6.22\n"
         "Longer contiguous signals are harder\nto explain by chance alignment gaps.",
         C_AMBER),
    ]

    for i, (term, rng, body, col) in enumerate(terms):
        left = Inches(0.4 + i * 4.2)
        add_rect(sl, left, Inches(2.05), Inches(3.95), Inches(0.48), col)
        add_text(sl, term, left + Inches(0.1), Inches(2.05), Inches(3.75), Inches(0.48),
                 font_size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, f"Range: {rng}", left + Inches(0.1), Inches(2.57), Inches(3.75), Inches(0.28),
                 font_size=Pt(11), color=col, bold=True)
        add_text(sl, body, left + Inches(0.1), Inches(2.88), Inches(3.75), Inches(1.5),
                 font_size=Pt(12), color=C_DARK)

    add_rect(sl, Inches(0.4), Inches(4.55), Inches(12.5), Inches(0.04), C_NAVY)

    tf = bullet_frame(sl, Inches(0.4), Inches(4.68), Inches(12.5), Inches(1.65),
                      RGBColor(0xEB, 0xF5, 0xFB))
    insights = [
        ("Why not use mean_divergence alone?", True, C_NAVY, Pt(14)),
        ("A repeat region or assembly gap also returns high divergence — not because it is pathogen-specific, "
         "but because it aligns poorly for technical reasons. The flank_conservation term corrects for this: "
         "if the surrounding 2 kb is shared between pathogen and commensal, the divergent block inside is "
         "structurally consistent with a genuine PAI insertion. This is exactly how HGT islands look — "
         "a foreign sequence inserted into an otherwise conserved genomic neighbourhood.", False, C_DARK, Pt(12)),
    ]
    first = True
    for text, bold, color, size in insights:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.size = size

    footer(sl, 8)
    return sl


def slide_09_thresholds(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_LIGHT)
    header_bar(sl, "Threshold Sensitivity Analysis",
               "Four threshold schemes — 500 bp and 1000 bp windows")

    add_text(sl,
             "500 bp windows",
             Inches(0.4), Inches(1.3), Inches(6.0), Inches(0.35),
             font_size=Pt(14), bold=True, color=C_NAVY)

    headers_500 = ["Scheme", "HIGH Windows", "% Genome", "Candidates", "Med. Length", "Med. Score"]
    rows_500 = [
        ["A (0.20/0.60)", "2,870", "25.65%", "143", "3,000 bp", "5.44"],
        ["B (0.15/0.50)", "3,160", "28.24%", "149", "2,500 bp", "5.35"],
        ["C (0.25/0.70)", "2,200", "19.66%", "121", "3,500 bp", "5.62"],
        ["D (tertiles)",  "3,720", "33.24%", "171", "2,000 bp", "5.18"],
    ]

    col_ws_500 = [1.5, 1.4, 1.1, 1.2, 1.4, 1.3]

    def draw_table(left, top, headers, rows, col_ws, highlight_row=0):
        t = top
        # header
        x = left
        for h, w in zip(headers, col_ws):
            add_rect(sl, x, t, Inches(w), Inches(0.34), C_NAVY)
            add_text(sl, h, x + Inches(0.05), t, Inches(w - 0.05), Inches(0.34),
                     font_size=Pt(10), bold=True, color=C_WHITE)
            x += Inches(w)
        t += Inches(0.34)
        for r_i, row in enumerate(rows):
            bg = RGBColor(0xD6, 0xEA, 0xF8) if r_i == highlight_row else C_WHITE
            x = left
            for cell, w in zip(row, col_ws):
                add_rect(sl, x, t, Inches(w), Inches(0.32), bg,
                         line_color=RGBColor(0xBD, 0xC3, 0xC7), line_width=Pt(0.5))
                add_text(sl, cell, x + Inches(0.05), t, Inches(w - 0.05), Inches(0.32),
                         font_size=Pt(10), color=C_DARK)
                x += Inches(w)
            t += Inches(0.32)

    draw_table(Inches(0.4), Inches(1.7), headers_500, rows_500, col_ws_500, highlight_row=0)

    add_text(sl, "1000 bp windows",
             Inches(0.4), Inches(4.0), Inches(6.0), Inches(0.35),
             font_size=Pt(14), bold=True, color=C_NAVY)

    rows_1000 = [
        ["A (0.20/0.60)", "2,620", "23.42%", "118", "4,000 bp", "5.81"],
        ["B (0.15/0.50)", "2,910", "26.01%", "125", "3,500 bp", "5.68"],
        ["C (0.25/0.70)", "2,050", "18.33%", "98",  "5,000 bp", "6.04"],
        ["D (tertiles)",  "3,400", "30.39%", "142", "3,000 bp", "5.52"],
    ]
    draw_table(Inches(0.4), Inches(4.4), headers_500, rows_1000, col_ws_500, highlight_row=0)

    # Right — interpretation
    tf = bullet_frame(sl, Inches(8.6), Inches(1.3), Inches(4.35), Inches(5.52))
    notes = [
        ("Interpretation", True, C_NAVY, Pt(14)),
        ("• Scheme A is the primary scheme (primary results).", False, C_DARK, Pt(12)),
        ("• Schemes B, C, D serve as sensitivity analysis —\n"
         "  verifying that conclusions hold across thresholds.", False, C_DARK, Pt(12)),
        ("• Scheme C (stricter HIGH threshold) yields fewer but\n"
         "  higher-confidence candidates with longer median length.", False, C_DARK, Pt(12)),
        ("• Scheme D (data-driven tertiles) has the highest\n"
         "  candidate count but lowest median marker score.", False, C_DARK, Pt(12)),
        ("", False, C_DARK, Pt(5)),
        ("Key constraint", True, C_RED, Pt(13)),
        ("Threshold choice is an empirical modelling decision.\n"
         "No single scheme is universally optimal — sensitivity\n"
         "analysis is mandatory before any threshold is cited\n"
         "in a thesis claim.", False, C_DARK, Pt(12)),
    ]
    first = True
    for text, bold, color, size in notes:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.size = size

    footer(sl, 9)
    return sl


def slide_10_landscape(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_LIGHT)
    header_bar(sl, "Chromosome Divergence Landscape",
               "Where HIGH-gradient regions fall on the Sakai genome")

    img = FIG_DIR / "divergence_landscape_chromosome.png"
    if img.exists():
        add_image(sl, img, Inches(0.3), Inches(1.25), Inches(9.5))

    tf = bullet_frame(sl, Inches(9.95), Inches(1.3), Inches(3.0), Inches(5.9))
    notes = [
        ("Reading the figure", True, C_NAVY, Pt(13)),
        ("Top panel: Mean divergence\nscore per 500 bp window.\nBlue=LOW, orange=MID,\nred=HIGH.", False, C_DARK, Pt(11)),
        ("", False, C_DARK, Pt(4)),
        ("Shaded bands: known PAI\ncoordinates from Hayashi\n2001 and Perna 2001.", False, C_DARK, Pt(11)),
        ("", False, C_DARK, Pt(4)),
        ("Middle strip: candidate\nregions (Scheme A) as\nfilled red bars.", False, C_DARK, Pt(11)),
        ("", False, C_DARK, Pt(4)),
        ("Bottom strip: pathogen-\nvs-pathogen negative\ncontrol divergence\n(expected lower).", False, C_DARK, Pt(11)),
        ("", False, C_DARK, Pt(4)),
        ("Observation", True, C_RED, Pt(13)),
        ("HIGH regions cluster\nnear annotated PAIs —\nespecially LEE and OI\nislands — confirming\nbiological specificity.", False, C_DARK, Pt(11)),
    ]
    first = True
    for text, bold, color, size in notes:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.size = size

    footer(sl, 10)
    return sl


def slide_11_tiers(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_LIGHT)
    header_bar(sl, "Confidence Tiers — Intersection Filter",
               "Combining NUCmer and minimap2 agreement into validation priority")

    img = FIG_DIR / "divergence_confidence_tiers.png"
    if img.exists():
        add_image(sl, img, Inches(0.3), Inches(1.25), Inches(9.0))

    tf = bullet_frame(sl, Inches(9.5), Inches(1.3), Inches(3.5), Inches(5.9))
    notes = [
        ("Tier system", True, C_NAVY, Pt(13)),
        ("Tier 1 (green): Both\npipelines agree.\n333 markers.\nHighest confidence.", False, C_GREEN, Pt(12)),
        ("", False, C_DARK, Pt(4)),
        ("Tier 2 (orange): NUCmer\ndetects divergence but\nminimap2 does not.\n82 markers.", False, C_AMBER, Pt(12)),
        ("", False, C_DARK, Pt(4)),
        ("Shapes:", True, C_NAVY, Pt(13)),
        ("◆ Diamond = DIVERGED\n  (positive class)\n● Circle = MODERATE\n  (negative class)", False, C_DARK, Pt(12)),
        ("", False, C_DARK, Pt(4)),
        ("Design decision", True, C_TEAL, Pt(13)),
        ("All 415 markers used\nin classifier. Tier 1\nused to PRIORITISE\nbiological validation,\nnot to filter ML input\n(no differential DIVERGED\nenrichment detected).", False, C_DARK, Pt(11)),
    ]
    first = True
    for text, bold, color, size in notes:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.size = size

    footer(sl, 11)
    return sl


def slide_12_ml(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_LIGHT)
    header_bar(sl, "ML Classifier — Features and Performance",
               "XGBoost with 15 features, StratifiedGroupKFold cross-validation")

    # Feature table
    add_text(sl, "15-Feature Vector",
             Inches(0.4), Inches(1.32), Inches(5.0), Inches(0.34),
             font_size=Pt(13), bold=True, color=C_NAVY)

    features = [
        ("NUCmer / Compositional (10)", [
            "blastn_identity", "cai_score", "gc_delta",
            "srna_density", "align_coverage", "kmer_deviation",
            "presence_pathogenic", "presence_non_pathogenic",
            "pangenome_score", "anvio_cluster_score",
        ], C_TEAL),
        ("Gradient Features (5)", [
            "mean_divergence", "flank_conservation_2000bp",
            "marker_score_2000bp", "proportion_high_windows",
            "proportion_mid_windows",
        ], C_AMBER),
    ]

    top_y = Inches(1.72)
    for group, feats, col in features:
        add_rect(sl, Inches(0.4), top_y, Inches(5.0), Inches(0.3), col)
        add_text(sl, group, Inches(0.5), top_y, Inches(4.8), Inches(0.3),
                 font_size=Pt(11), bold=True, color=C_WHITE)
        top_y += Inches(0.3)
        for f in feats:
            add_rect(sl, Inches(0.4), top_y, Inches(5.0), Inches(0.28),
                     C_WHITE, line_color=RGBColor(0xBD, 0xC3, 0xC7), line_width=Pt(0.5))
            add_text(sl, f, Inches(0.55), top_y, Inches(4.7), Inches(0.28),
                     font_size=Pt(10), color=C_DARK)
            top_y += Inches(0.28)
        top_y += Inches(0.1)

    # AUROC callouts
    for left, val, label, col in [
        (Inches(5.8),  "0.717", "NUCmer baseline\n(10 features)", C_TEAL),
        (Inches(8.55), "0.709", "Combined model\n(15 features)", C_AMBER),
        (Inches(11.3), "0.549", "Gradient-only\n(5 features)", C_RED),
    ]:
        add_rect(sl, left, Inches(1.35), Inches(2.4), Inches(1.8), col)
        add_text(sl, val, left, Inches(1.4), Inches(2.4), Inches(0.95),
                 font_size=Pt(36), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, "AUROC", left, Inches(2.3), Inches(2.4), Inches(0.3),
                 font_size=Pt(11), color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, label, left, Inches(2.62), Inches(2.4), Inches(0.5),
                 font_size=Pt(11), color=C_WHITE, align=PP_ALIGN.CENTER)

    img = FIG_DIR / "roc_comparison.png"
    if img.exists():
        add_image(sl, img, Inches(5.8), Inches(3.3), Inches(7.15), Inches(3.4))

    footer(sl, 12)
    return sl


def slide_13_ablation(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_LIGHT)
    header_bar(sl, "Ablation Study — Feature Contributions",
               "Systematically removing feature subsets to measure individual contribution")

    add_text(sl,
             "An ablation study removes one group of features at a time and re-trains the classifier. "
             "If performance drops sharply, that group was carrying important signal.",
             Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.55),
             font_size=Pt(13), color=C_DARK)

    # Ablation results bars
    ablation = [
        ("NUCmer 10-feature baseline", 0.717, C_TEAL,  "Best single pipeline"),
        ("Combined 15-feature model",  0.709, C_AMBER, "Within CV uncertainty of baseline"),
        ("Gradient-only (5 features)", 0.549, C_RED,   "Signal present but weaker"),
        ("Label shuffle (null)",        0.428, RGBColor(0x95, 0xA5, 0xA6), "Near-random — no leakage"),
    ]

    bar_left = Inches(0.4)
    bar_top  = Inches(2.0)
    bar_h    = Inches(0.72)
    bar_gap  = Inches(0.22)
    max_val  = 1.0
    bar_max_w = Inches(8.0)

    for i, (label, auroc, col, note) in enumerate(ablation):
        t = bar_top + i * (bar_h + bar_gap)
        bar_w = Emu(int(bar_max_w * auroc / max_val))
        add_rect(sl, bar_left, t, bar_w, bar_h, col)
        add_text(sl, f"AUROC {auroc:.3f}",
                 bar_left + bar_w + Inches(0.12), t + Inches(0.08),
                 Inches(1.5), Inches(0.4),
                 font_size=Pt(15), bold=True, color=col)
        add_text(sl, label,
                 bar_left, t - Inches(0.3), Inches(8.5), Inches(0.28),
                 font_size=Pt(12), bold=True, color=C_DARK)
        add_text(sl, note,
                 bar_left + bar_w + Inches(0.12), t + Inches(0.38),
                 Inches(4.0), Inches(0.28),
                 font_size=Pt(11), color=C_DARK)

    img = FIG_DIR / "shap_comparison_core10_vs_combined15.png"
    if img.exists():
        add_image(sl, img, Inches(8.7), Inches(1.9), Inches(4.25), Inches(4.5))
        add_text(sl, "SHAP feature importance:\nNUCmer baseline (blue) vs\ncombined model (orange)",
                 Inches(8.7), Inches(6.42), Inches(4.25), Inches(0.45),
                 font_size=Pt(10), color=C_DARK, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(0.4), Inches(6.35), Inches(8.1), Inches(0.48),
             RGBColor(0xEB, 0xF5, 0xFB))
    add_text(sl,
             "Key finding: gradient features add no net improvement over the NUCmer baseline. "
             "Their standalone signal (0.549 vs null 0.428) confirms biological content — "
             "but that content is largely already captured by compositional features.",
             Inches(0.5), Inches(6.38), Inches(7.9), Inches(0.42),
             font_size=Pt(11), color=C_NAVY)

    footer(sl, 13)
    return sl


def slide_14_neg_control(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_LIGHT)
    header_bar(sl, "Negative Control — Specificity Validation",
               "Confirming that HIGH regions are pathogen-specific, not alignment artefacts")

    for left, val, label, sub, col in [
        (Inches(0.4),  "25.65%", "Pathogen vs Commensal",
         "Sakai aligned against 30 commensal strains\n(primary experiment)", C_RED),
        (Inches(4.6),  "2.35%",  "Pathogen vs Pathogen",
         "Sakai aligned against 5 EHEC/pathogenic strains\n(negative control)", C_TEAL),
        (Inches(8.8),  "10.9×",  "Enrichment ratio",
         "Fold-increase in HIGH windows:\ncommensal comparison over pathogen comparison", C_AMBER),
    ]:
        add_rect(sl, left, Inches(1.3), Inches(3.85), Inches(2.2), col)
        add_text(sl, val, left, Inches(1.35), Inches(3.85), Inches(1.0),
                 font_size=Pt(40), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, "HIGH windows", left, Inches(2.28), Inches(3.85), Inches(0.3),
                 font_size=Pt(11), color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, label, left, Inches(2.6), Inches(3.85), Inches(0.35),
                 font_size=Pt(12), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, sub, left + Inches(0.1), Inches(2.98), Inches(3.65), Inches(0.5),
                 font_size=Pt(10), color=C_WHITE, align=PP_ALIGN.CENTER)

    tf = bullet_frame(sl, Inches(0.4), Inches(3.72), Inches(12.5), Inches(3.1))
    notes = [
        ("What this means", True, C_NAVY, Pt(14)),
        ("The HIGH-gradient regions are not simply \"hard to align\" — they are specifically absent from "
         "commensal strains. When Sakai is aligned against other EHEC pathogenic strains, only 2.35% of "
         "windows are HIGH, because those strains share the same PAIs. When aligned against commensals, "
         "25.65% are HIGH — driven by the genuine absence of pathogen-specific genomic islands from "
         "commensal genomes. The 10.9× enrichment ratio confirms the pipeline is detecting pathogen-"
         "specific divergence, not technical noise.", False, C_DARK, Pt(13)),
        ("", False, C_DARK, Pt(6)),
        ("Label shuffle control", True, C_TEAL, Pt(14)),
        ("Randomly permuting ML labels produced AUROC 0.428 — close to the 0.50 expected for a null "
         "model. This confirms that the classifier's AUROC 0.709 is driven by real biological signal, "
         "not by label imbalance, spatial autocorrelation, or CV leakage.", False, C_DARK, Pt(13)),
    ]
    first = True
    for text, bold, color, size in notes:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.size = size

    footer(sl, 14)
    return sl


def slide_15_conclusion(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C_NAVY)

    add_text(sl, "Conclusions and Future Directions",
             Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.65),
             font_size=Pt(26), bold=True, color=C_WHITE)
    add_rect(sl, Inches(0.5), Inches(0.98), Inches(12.3), Inches(0.04), C_TEAL)

    conclusions = [
        ("What was achieved", C_TEAL, [
            "A minimap2 divergence gradient pipeline extracting 3-state (LOW/MID/HIGH) "
            "window-level features from 30 commensal comparisons",
            "415 candidate marker regions ranked by a composite marker_score "
            "(divergence × flank conservation × log-length)",
            "Negative control confirming 10.9× enrichment of HIGH windows in "
            "pathogen-vs-commensal vs pathogen-vs-pathogen",
            "Gradient features confirmed biologically informative (AUROC 0.549 standalone)",
        ]),
        ("What was found", C_AMBER, [
            "Gradient features do not improve over NUCmer 10-feature baseline "
            "(0.709 combined vs 0.717 NUCmer — within CV uncertainty)",
            "333 Tier-1 markers where both pipelines agree: highest-confidence validation targets",
            "Threshold choice is empirical — all 4 schemes show consistent enrichment pattern",
        ]),
        ("Future directions", C_GREEN, [
            "Experimental validation: PCR/WGS confirmation of top Tier-1 marker presence/absence",
            "Extend panel: test on non-E. coli pathogens to assess generalisability",
            "Integrate structural variation signals (SVs, inversions) beyond single-base divergence",
        ]),
    ]

    top_y = Inches(1.1)
    for title, col, points in conclusions:
        add_rect(sl, Inches(0.5), top_y, Inches(12.3), Inches(0.36), col)
        add_text(sl, title, Inches(0.65), top_y, Inches(12.0), Inches(0.36),
                 font_size=Pt(13), bold=True, color=C_WHITE)
        top_y += Inches(0.36)
        for pt in points:
            add_rect(sl, Inches(0.5), top_y, Inches(12.3), Inches(0.36),
                     RGBColor(0x2C, 0x3E, 0x50))
            add_text(sl, f"• {pt}", Inches(0.65), top_y, Inches(12.0), Inches(0.36),
                     font_size=Pt(11), color=C_WHITE)
            top_y += Inches(0.36)
        top_y += Inches(0.12)

    add_text(sl,
             "The principal contribution of this work is not a new biological concept but a new operational "
             "representation: the three-state divergence gradient as a structured, interpretable, and "
             "mechanistically grounded feature for pathogen-marker discovery.",
             Inches(0.5), top_y + Inches(0.05), Inches(12.3), Inches(0.55),
             font_size=Pt(12), color=C_TEAL, align=PP_ALIGN.CENTER)

    footer(sl, 15)
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    builders = [
        slide_01_title,
        slide_02_problem,
        slide_03_objectives,
        slide_04_biology,
        slide_05_gap,
        slide_06_pipeline,
        slide_07_gradient,
        slide_08_marker_score,
        slide_09_thresholds,
        slide_10_landscape,
        slide_11_tiers,
        slide_12_ml,
        slide_13_ablation,
        slide_14_neg_control,
        slide_15_conclusion,
    ]

    for i, builder in enumerate(builders, 1):
        print(f"  Building slide {i:02d}: {builder.__name__[6:]}")
        builder(prs)

    prs.save(str(OUT_PATH))
    print(f"\nSaved: {OUT_PATH}  ({OUT_PATH.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    print("Generating thesis presentation (15 slides)...")
    main()
